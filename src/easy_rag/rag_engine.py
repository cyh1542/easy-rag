from __future__ import annotations

import json
import math
import shutil
import re
import xml.etree.ElementTree as ET
from concurrent.futures import ThreadPoolExecutor, as_completed
from pathlib import Path
from typing import Any

import chromadb
import pandas as pd
import pymysql
from bs4 import BeautifulSoup
from huggingface_hub import snapshot_download
from openai import OpenAI
from pymysql.cursors import DictCursor
from pypdf import PdfReader
from docx import Document
from pptx import Presentation

from easy_rag.config import Settings, apply_hf_hub_settings
from easy_rag.knowledge_bases import (
    ChunkConfig,
    KnowledgeBaseProfile,
    chunk_config_from_settings,
    merge_retrieval_results,
    resolve_retrieval_profiles,
)
from easy_rag.logger_config import setup_logging
from easy_rag.timing_utils import StageTimer


SUPPORTED_EXTENSIONS = {
    ".txt",
    ".md",
    ".pdf",
    ".csv",
    ".xlsx",
    ".xls",
    ".xlsm",
    ".docx",
    ".pptx",
    ".html",
    ".json",
    ".xml",
}
MYSQL_IDENTIFIER_PATTERN = re.compile(r"^[A-Za-z0-9_]+$")
STRUCTURE_SPLIT_PATTERN = re.compile(
    r"(?=^#{1,6}\s+|^记录 \d+|^第[0-9一二三四五六七八九十百千]+[章节部分条]|^\d+\.\s+|^[一二三四五六七八九十]+、)",
    re.MULTILINE,
)
RECURSIVE_SEPARATORS = ["\n\n", "\n", "。", "！", "？", ".", "!", "?", "；", ";", " ", ""]
FILE_PATH_METADATA_KEYS = (
    "file_name",
    "file_stem",
    "file_suffix",
    "relative_path",
    "parent_dir",
    "document_type",
)
logger = setup_logging()


def _build_file_path_metadata(file_path: Path, knowledge_dir: Path) -> dict[str, str]:
    resolved = file_path.resolve()
    knowledge_root = knowledge_dir.resolve()
    try:
        relative_path = resolved.relative_to(knowledge_root).as_posix()
    except ValueError:
        relative_path = resolved.as_posix()

    parent = Path(relative_path).parent
    parent_dir = "" if str(parent) in {"", "."} else parent.as_posix()

    return {
        "source": file_path.name,
        "path": str(resolved),
        "file_name": file_path.name,
        "file_stem": file_path.stem,
        "file_suffix": file_path.suffix.lower(),
        "relative_path": relative_path,
        "parent_dir": parent_dir,
        "document_type": "file",
    }


def _build_mysql_document_metadata(
    *,
    source: str,
    path: str,
    document_type: str,
) -> dict[str, str]:
    return {
        "source": source,
        "path": path,
        "file_name": source,
        "file_stem": source,
        "file_suffix": "",
        "relative_path": source,
        "parent_dir": "",
        "document_type": document_type,
    }


def _chunk_metadata_from_document(
    document: dict[str, Any],
    *,
    chunk_index: int,
    kb_id: str = "",
    collection_name: str = "",
) -> dict[str, Any]:
    metadata: dict[str, Any] = {
        "source": document.get("source", ""),
        "path": document.get("path", ""),
        "chunk_index": chunk_index,
        "kb_id": kb_id,
        "collection_name": collection_name,
    }
    for key in FILE_PATH_METADATA_KEYS:
        value = document.get(key)
        if value not in (None, ""):
            metadata[key] = value
    return metadata


def _context_item_from_metadata(
    metadata: dict[str, Any],
    content: str,
    **extra: Any,
) -> dict[str, Any]:
    item: dict[str, Any] = {
        "content": content,
        "source": str(metadata.get("source", "unknown")),
        "path": str(metadata.get("path", "")),
    }
    for key in FILE_PATH_METADATA_KEYS:
        value = metadata.get(key)
        if value not in (None, ""):
            item[key] = value
    chunk_index = metadata.get("chunk_index")
    if chunk_index is not None:
        item["chunk_index"] = chunk_index
    item.update(extra)
    return item


def _format_context_label(item: dict[str, Any]) -> str:
    return str(item.get("relative_path") or item.get("file_name") or item.get("source") or "unknown")


def _seed_fused_item(item: dict[str, Any]) -> dict[str, Any]:
    seeded: dict[str, Any] = {
        "content": item.get("content", ""),
        "source": item.get("source", "unknown"),
        "path": item.get("path", ""),
        "score": 0.0,
    }
    for key in FILE_PATH_METADATA_KEYS:
        value = item.get(key)
        if value not in (None, ""):
            seeded[key] = value
    chunk_index = item.get("chunk_index")
    if chunk_index is not None:
        seeded["chunk_index"] = chunk_index
    return seeded

DEFAULT_CHAT_SYSTEM_PROMPT = (
    "你是一个严谨的中文知识库问答助手。"
    "请优先依据检索到的上下文回答问题。"
    "如果上下文不足以支持结论，请明确说明不知道或信息不足，"
    "不要编造事实。"
)


class EasyRAG:
    def __init__(self, settings: Settings) -> None:
        self.settings = settings
        self._validate_settings()
        apply_hf_hub_settings(self.settings)
        self.settings.knowledge_dir.mkdir(parents=True, exist_ok=True)
        self.settings.chroma_dir.mkdir(parents=True, exist_ok=True)

        self._openai_client: OpenAI | None = None
        self._embedding_client: OpenAI | None = None
        self._local_embedding_model: Any | None = None
        self._vector_client: Any | None = None
        self._active_chunk_config: ChunkConfig | None = None

    def _chunk(self) -> ChunkConfig:
        if self._active_chunk_config is not None:
            return self._active_chunk_config
        return chunk_config_from_settings(self.settings)

    def _validate_settings(self) -> None:
        if self.settings.chunk_overlap >= self.settings.chunk_size:
            raise ValueError("CHUNK_OVERLAP 必须小于 CHUNK_SIZE。")

        if self.settings.chunk_strategy not in {"fixed", "recursive", "semantic", "structure"}:
            raise ValueError("CHUNK_STRATEGY 只支持 fixed、recursive、semantic 或 structure。")

        if not 0 < self.settings.semantic_chunk_threshold <= 1:
            raise ValueError("SEMANTIC_CHUNK_THRESHOLD 必须在 0 到 1 之间。")

        if self.settings.embedding_provider not in {"remote", "local", "huggingface"}:
            raise ValueError("EMBEDDING_PROVIDER 只支持 remote、local 或 huggingface。")

        if self.settings.embedding_provider == "remote" and not self.settings.embedding_model:
            raise ValueError("当 EMBEDDING_PROVIDER=remote 时，必须填写 EMBEDDING_MODEL。")

        if self.settings.embedding_provider == "local" and not self.settings.local_embedding_model:
            raise ValueError("当 EMBEDDING_PROVIDER=local 时，必须填写 LOCAL_EMBEDDING_MODEL。")

        if self.settings.embedding_provider == "huggingface" and not self.settings.hf_embedding_repo_id:
            raise ValueError("当 EMBEDDING_PROVIDER=huggingface 时，必须填写 HF_EMBEDDING_REPO_ID。")

        if self.settings.retrieval_method not in {"keyword", "vector", "rerank", "rrf"}:
            raise ValueError("RETRIEVAL_METHOD 只支持 keyword、vector、rerank 或 rrf。")

        if self.settings.rerank_candidate_k < self.settings.top_k:
            raise ValueError("RERANK_CANDIDATE_K 必须大于或等于 TOP_K。")

        if self.settings.rrf_k <= 0:
            raise ValueError("RRF_K 必须大于 0。")

        if self.settings.mysql_enabled:
            missing_fields: list[str] = []
            if not self.settings.mysql_host:
                missing_fields.append("MYSQL_HOST")
            if not self.settings.mysql_user:
                missing_fields.append("MYSQL_USER")
            if not self.settings.mysql_database:
                missing_fields.append("MYSQL_DATABASE")
            if missing_fields:
                raise ValueError(
                    "启用 MYSQL_ENABLED 后请补全这些配置: " + ", ".join(missing_fields)
                )
            if not self.settings.mysql_tables and not self.settings.mysql_query:
                raise ValueError("启用 MYSQL_ENABLED 后，请至少配置 MYSQL_TABLES 或 MYSQL_QUERY。")

    def _validate_openai_api_key(self) -> None:
        if not self.settings.api_key or self.settings.api_key == "your_api_key_here":
            raise ValueError("请先在配置中填写有效的 OPENAI_API_KEY。")

    def _validate_embedding_api_key(self) -> None:
        if self.settings.embedding_provider != "remote":
            return

        configured_key = self.settings.embedding_api_key.strip()
        if configured_key:
            if configured_key == "your_api_key_here":
                raise ValueError("请先在配置中填写有效的 EMBEDDING_API_KEY。")
            return

        if not self.settings.api_key or self.settings.api_key == "your_api_key_here":
            raise ValueError("请先在配置中填写有效的 EMBEDDING_API_KEY。")

    def _resolve_embedding_api_key(self) -> str:
        self._validate_embedding_api_key()
        configured_key = self.settings.embedding_api_key.strip()
        if configured_key:
            return configured_key
        return self.settings.api_key

    def _resolve_chat_system_prompt(self) -> str:
        custom = self.settings.chat_thinking_prompt.strip()
        return custom or DEFAULT_CHAT_SYSTEM_PROMPT

    def _format_context_block(self, index: int, item: dict[str, Any]) -> str:
        label = _format_context_label(item)
        lines = [f"[片段 {index}]", f"来源: {label}"]
        file_path = str(item.get("path", "")).strip()
        if file_path and file_path != label:
            lines.append(f"路径: {file_path}")
        parent_dir = str(item.get("parent_dir", "")).strip()
        if parent_dir:
            lines.append(f"目录: {parent_dir}")
        lines.append(f"内容: {item.get('content', '')}")
        return "\n".join(lines)

    def _get_openai_client(self) -> OpenAI:
        self._validate_openai_api_key()

        if self._openai_client is None:
            self._openai_client = OpenAI(
                api_key=self.settings.api_key,
                base_url=self.settings.base_url,
            )
        return self._openai_client

    def _get_embedding_client(self) -> OpenAI:
        embedding_api_key = self._resolve_embedding_api_key()

        if self._embedding_client is None:
            self._embedding_client = OpenAI(
                api_key=embedding_api_key,
                base_url=self.settings.embedding_base_url or self.settings.base_url,
            )
        return self._embedding_client

    def _get_local_embedding_model(self) -> Any:
        if self._local_embedding_model is None:
            try:
                from sentence_transformers import SentenceTransformer
            except ImportError as exc:
                raise ImportError(
                    "未安装 sentence-transformers，请先执行 pip install -r requirements.txt。"
                ) from exc

            if self.settings.hf_endpoint:
                logger.info("Hugging Face Hub 镜像: %s", self.settings.hf_endpoint)
            logger.info(
                "本地 Embedding 模型加载中: model=%s device=%s",
                self.settings.local_embedding_model,
                self.settings.local_embedding_device or "auto",
            )
            self._local_embedding_model = SentenceTransformer(
                self.settings.local_embedding_model,
                device=self.settings.local_embedding_device or None,
            )
        return self._local_embedding_model

    def _download_huggingface_model(self) -> str:
        self.settings.hf_embedding_cache_dir.mkdir(parents=True, exist_ok=True)
        if self.settings.hf_endpoint:
            logger.info("Hugging Face Hub 镜像: %s", self.settings.hf_endpoint)
        logger.info("开始下载 Hugging Face embedding 模型: repo=%s", self.settings.hf_embedding_repo_id)
        return snapshot_download(
            repo_id=self.settings.hf_embedding_repo_id,
            local_dir=str(self.settings.hf_embedding_cache_dir / self.settings.hf_embedding_repo_id.replace("/", "__")),
            local_dir_use_symlinks=False,
        )

    def _get_huggingface_embedding_model(self) -> Any:
        if self._local_embedding_model is None:
            try:
                from sentence_transformers import SentenceTransformer
            except ImportError as exc:
                raise ImportError(
                    "未安装 sentence-transformers，请先执行 pip install -r requirements.txt。"
                ) from exc

            downloaded_path = self._download_huggingface_model()
            logger.info("Hugging Face Embedding 模型加载中: path=%s", downloaded_path)
            self._local_embedding_model = SentenceTransformer(
                downloaded_path,
                device=self.settings.local_embedding_device or None,
            )
        return self._local_embedding_model

    def list_supported_files(self, directory: Path | None = None) -> list[Path]:
        target_dir = directory or self.settings.knowledge_dir
        if not target_dir.exists():
            return []

        files: list[Path] = []
        for file_path in sorted(target_dir.rglob("*")):
            if file_path.is_file() and file_path.suffix.lower() in SUPPORTED_EXTENSIONS:
                files.append(file_path)
        return files

    def _read_text_file(self, file_path: Path) -> str:
        for encoding in ("utf-8", "utf-8-sig", "gbk", "gb18030"):
            try:
                return file_path.read_text(encoding=encoding)
            except UnicodeDecodeError:
                continue
        return file_path.read_text(encoding="utf-8", errors="ignore")

    def _read_pdf_file(self, file_path: Path) -> str:
        reader = PdfReader(str(file_path))
        pages = [(page.extract_text() or "").strip() for page in reader.pages]
        return "\n".join(page for page in pages if page)

    def _dataframe_to_text(self, dataframe: pd.DataFrame, title: str | None = None) -> str:
        if dataframe.empty and len(dataframe.columns) == 0:
            return ""

        normalized = dataframe.fillna("").astype(str)
        normalized.columns = [str(column) for column in normalized.columns]
        table_text = normalized.to_csv(index=False).strip() or "(空表)"

        if title:
            return f"{title}\n{table_text}"
        return table_text

    def _rows_to_text(self, rows: list[dict[str, Any]], title: str | None = None) -> str:
        if not rows:
            return ""

        sections: list[str] = []
        for index, row in enumerate(rows, start=1):
            row_lines = [f"{column}: {value if value is not None else ''}" for column, value in row.items()]
            sections.append(f"记录 {index}\n" + "\n".join(row_lines))

        content = "\n\n".join(item.strip() for item in sections if item.strip())
        if title:
            return f"{title}\n\n{content}".strip()
        return content

    def _read_csv_file(self, file_path: Path) -> str:
        last_error: Exception | None = None
        for encoding in ("utf-8", "utf-8-sig", "gbk", "gb18030"):
            try:
                dataframe = pd.read_csv(file_path, encoding=encoding)
                return self._dataframe_to_text(dataframe)
            except UnicodeDecodeError as exc:
                last_error = exc
                continue

        if last_error:
            raise last_error
        raise ValueError(f"无法读取 CSV 文件: {file_path}")

    def _read_excel_file(self, file_path: Path) -> str:
        workbook = pd.read_excel(file_path, sheet_name=None)
        sections: list[str] = []

        for sheet_name, dataframe in workbook.items():
            section = self._dataframe_to_text(dataframe, title=f"工作表: {sheet_name}")
            if section:
                sections.append(section)

        return "\n\n".join(sections)

    def _read_docx_file(self, file_path: Path) -> str:
        document = Document(str(file_path))
        paragraphs = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
        return "\n\n".join(paragraphs)

    def _read_pptx_file(self, file_path: Path) -> str:
        presentation = Presentation(str(file_path))
        slides_text: list[str] = []

        for slide_index, slide in enumerate(presentation.slides, start=1):
            texts: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    content = shape.text.strip()
                    if content:
                        texts.append(content)
            if texts:
                slides_text.append(f"幻灯片 {slide_index}\n" + "\n".join(texts))

        return "\n\n".join(slides_text)

    def _read_html_file(self, file_path: Path) -> str:
        html = self._read_text_file(file_path)
        soup = BeautifulSoup(html, "html.parser")
        text = soup.get_text(separator="\n", strip=True)
        lines = [line.strip() for line in text.splitlines() if line.strip()]
        return "\n".join(lines)

    def _read_json_file(self, file_path: Path) -> str:
        raw_text = self._read_text_file(file_path)
        data = json.loads(raw_text)
        return json.dumps(data, ensure_ascii=False, indent=2)

    def _flatten_xml_element(self, element: ET.Element, depth: int = 0) -> list[str]:
        indent = "  " * depth
        lines = [f"{indent}<{element.tag}>"]

        text = (element.text or "").strip()
        if text:
            lines.append(f"{indent}{text}")

        for child in list(element):
            lines.extend(self._flatten_xml_element(child, depth + 1))

        return lines

    def _read_xml_file(self, file_path: Path) -> str:
        root = ET.parse(file_path).getroot()
        return "\n".join(self._flatten_xml_element(root))

    def _read_file(self, file_path: Path) -> str:
        suffix = file_path.suffix.lower()
        if suffix in {".txt", ".md"}:
            return self._read_text_file(file_path)
        if suffix == ".pdf":
            return self._read_pdf_file(file_path)
        if suffix == ".csv":
            return self._read_csv_file(file_path)
        if suffix in {".xlsx", ".xls", ".xlsm"}:
            return self._read_excel_file(file_path)
        if suffix == ".docx":
            return self._read_docx_file(file_path)
        if suffix == ".pptx":
            return self._read_pptx_file(file_path)
        if suffix == ".html":
            return self._read_html_file(file_path)
        if suffix == ".json":
            return self._read_json_file(file_path)
        if suffix == ".xml":
            return self._read_xml_file(file_path)
        raise ValueError(f"暂不支持的文件类型: {file_path.suffix}")

    def _safe_mysql_identifier(self, name: str) -> str:
        if not MYSQL_IDENTIFIER_PATTERN.fullmatch(name):
            raise ValueError(f"MySQL 表名只支持字母、数字和下划线，当前值不安全: {name}")
        return name

    def _get_mysql_connection(self) -> pymysql.connections.Connection:
        return pymysql.connect(
            host=self.settings.mysql_host,
            port=self.settings.mysql_port,
            user=self.settings.mysql_user,
            password=self.settings.mysql_password,
            database=self.settings.mysql_database,
            charset=self.settings.mysql_charset,
            cursorclass=DictCursor,
        )

    def test_mysql_connection(self) -> dict[str, Any]:
        if not self.settings.mysql_enabled:
            raise ValueError("请先启用 MYSQL_ENABLED。")

        logger.info("开始测试 MySQL 连接: host=%s database=%s", self.settings.mysql_host, self.settings.mysql_database)
        with self._get_mysql_connection() as connection:
            with connection.cursor() as cursor:
                cursor.execute("SELECT DATABASE() AS database_name, VERSION() AS version")
                info = cursor.fetchone() or {}
                cursor.execute("SHOW TABLES")
                raw_tables = cursor.fetchall()

        table_names = [str(next(iter(item.values()))) for item in raw_tables[:20]]
        return {
            "database": info.get("database_name", self.settings.mysql_database),
            "version": info.get("version", "unknown"),
            "table_count": len(raw_tables),
            "sample_tables": table_names,
        }

    def preview_mysql_table(self, table_name: str, limit: int = 10) -> dict[str, Any]:
        if not self.settings.mysql_enabled:
            raise ValueError("请先启用 MYSQL_ENABLED。")

        safe_table_name = self._safe_mysql_identifier(table_name)
        logger.info(
            "开始预览 MySQL 表数据: table=%s limit=%s",
            safe_table_name,
            limit,
        )

        sql = f"SELECT * FROM `{safe_table_name}` LIMIT %s"
        with self._get_mysql_connection() as connection:
            with connection.cursor() as cursor:
                cursor.execute(sql, (limit,))
                rows = list(cursor.fetchall())

        columns = list(rows[0].keys()) if rows else []
        return {
            "table_name": safe_table_name,
            "columns": columns,
            "rows": rows,
            "row_count": len(rows),
        }

    def load_file_documents(self) -> list[dict[str, str]]:
        documents: list[dict[str, str]] = []
        for file_path in self.list_supported_files():
            content = self._read_file(file_path).strip()
            if not content:
                continue

            documents.append(
                {
                    **_build_file_path_metadata(file_path, self.settings.knowledge_dir),
                    "content": content.replace("\r\n", "\n").replace("\r", "\n"),
                }
            )
        return documents

    def load_mysql_documents(self) -> list[dict[str, str]]:
        documents: list[dict[str, str]] = []
        if not self.settings.mysql_enabled:
            return documents

        with self._get_mysql_connection() as connection:
            if self.settings.mysql_query:
                with connection.cursor() as cursor:
                    cursor.execute(self.settings.mysql_query)
                    rows = list(cursor.fetchall())

                content = self._rows_to_text(rows, title="MySQL 自定义查询结果")
                if content:
                    documents.append(
                        {
                            **_build_mysql_document_metadata(
                                source="mysql_query",
                                path=(
                                    f"mysql://{self.settings.mysql_host}:{self.settings.mysql_port}/"
                                    f"{self.settings.mysql_database}/query"
                                ),
                                document_type="mysql_query",
                            ),
                            "content": content,
                        }
                    )

            for raw_table_name in self.settings.mysql_tables:
                table_name = self._safe_mysql_identifier(raw_table_name)
                sql = f"SELECT * FROM `{table_name}` LIMIT %s"
                with connection.cursor() as cursor:
                    cursor.execute(sql, (self.settings.mysql_limit_per_table,))
                    rows = list(cursor.fetchall())

                content = self._rows_to_text(rows, title=f"MySQL 表: {table_name}")
                if content:
                    documents.append(
                        {
                            **_build_mysql_document_metadata(
                                source=f"mysql_table:{table_name}",
                                path=(
                                    f"mysql://{self.settings.mysql_host}:{self.settings.mysql_port}/"
                                    f"{self.settings.mysql_database}/{table_name}"
                                ),
                                document_type="mysql_table",
                            ),
                            "content": content,
                        }
                    )

        return documents

    def load_documents(self) -> tuple[list[dict[str, str]], dict[str, int]]:
        file_documents = self.load_file_documents()
        mysql_documents = self.load_mysql_documents()
        documents = file_documents + mysql_documents
        return documents, {
            "file_documents": len(file_documents),
            "mysql_documents": len(mysql_documents),
            "documents": len(documents),
        }

    def split_text(self, text: str, chunk_config: ChunkConfig | None = None) -> list[str]:
        normalized = text.replace("\r\n", "\n").replace("\r", "\n").strip()
        if not normalized:
            return []

        previous = self._active_chunk_config
        self._active_chunk_config = chunk_config or chunk_config_from_settings(self.settings)
        try:
            strategy = self._chunk().chunk_strategy
            logger.debug("当前切分策略: %s", strategy)
            if strategy == "fixed":
                return self._split_text_fixed(normalized)
            if strategy == "recursive":
                return self._split_text_recursive(normalized)
            if strategy == "semantic":
                return self._split_text_semantic(normalized)
            return self._split_text_structure(normalized)
        finally:
            self._active_chunk_config = previous

    def _split_text_fixed(self, text: str) -> list[str]:
        paragraphs = [item.strip() for item in text.split("\n\n") if item.strip()]
        if not paragraphs:
            return []

        chunks: list[str] = []
        current = ""

        for paragraph in paragraphs:
            if len(paragraph) <= self._chunk().chunk_size:
                candidate = f"{current}\n\n{paragraph}".strip() if current else paragraph
                if len(candidate) <= self._chunk().chunk_size:
                    current = candidate
                    continue

                if current:
                    chunks.append(current)
                current = paragraph
                continue

            if current:
                chunks.append(current)
                current = ""

            start = 0
            while start < len(paragraph):
                end = min(start + self._chunk().chunk_size, len(paragraph))
                piece = paragraph[start:end].strip()
                if piece:
                    chunks.append(piece)
                if end >= len(paragraph):
                    break
                start = end - self._chunk().chunk_overlap

        if current:
            chunks.append(current)

        return chunks

    def _apply_chunk_overlap(self, chunks: list[str]) -> list[str]:
        if self._chunk().chunk_overlap <= 0 or len(chunks) <= 1:
            return chunks

        overlapped = [chunks[0]]
        for index in range(1, len(chunks)):
            prefix = chunks[index - 1][-self._chunk().chunk_overlap :]
            piece = chunks[index]
            overlapped.append(f"{prefix}{piece}" if prefix else piece)
        return overlapped

    def _recursive_split(self, text: str, separators: list[str]) -> list[str]:
        cleaned = text.strip()
        if not cleaned:
            return []
        if len(cleaned) <= self._chunk().chunk_size:
            return [cleaned]

        separator = separators[-1]
        next_separators = separators[-1:]
        for index, candidate in enumerate(separators):
            if candidate in cleaned:
                separator = candidate
                next_separators = separators[index + 1 :] or [""]
                break

        if separator == "":
            return self._split_text_fixed(cleaned)

        splits = cleaned.split(separator)
        merged: list[str] = []
        current = ""
        for item in splits:
            segment = item.strip()
            if not segment:
                continue
            joined = f"{current}{separator}{segment}".strip() if current else segment
            if len(joined) <= self._chunk().chunk_size:
                current = joined
                continue
            if current:
                merged.append(current)
            current = segment
        if current:
            merged.append(current)

        final_chunks: list[str] = []
        for piece in merged:
            if len(piece) <= self._chunk().chunk_size:
                final_chunks.append(piece)
            else:
                final_chunks.extend(self._recursive_split(piece, next_separators))
        return final_chunks

    def _split_text_recursive(self, text: str) -> list[str]:
        chunks = self._recursive_split(text, RECURSIVE_SEPARATORS)
        return self._apply_chunk_overlap(chunks)

    def _split_into_sentences(self, text: str) -> list[str]:
        parts = re.split(r"(?<=[。！？!?；;\n])", text)
        return [part.strip() for part in parts if part.strip()]

    def _split_text_semantic(self, text: str) -> list[str]:
        sentences = self._split_into_sentences(text)
        if not sentences:
            return []
        if len(sentences) == 1:
            if len(sentences[0]) <= self._chunk().chunk_size:
                return sentences
            return self._split_text_fixed(sentences[0])

        embeddings = self._embed_texts(sentences)
        groups: list[str] = []
        current = sentences[0]
        current_embedding = embeddings[0]

        for index in range(1, len(sentences)):
            similarity = self._cosine_similarity(current_embedding, embeddings[index])
            candidate = f"{current}{sentences[index]}"
            if (
                similarity >= self._chunk().semantic_chunk_threshold
                and len(candidate) <= self._chunk().chunk_size
            ):
                current = candidate
                current_embedding = embeddings[index]
                continue

            if current.strip():
                groups.append(current.strip())
            current = sentences[index]
            current_embedding = embeddings[index]

        if current.strip():
            groups.append(current.strip())

        final_chunks: list[str] = []
        for group in groups:
            if len(group) <= self._chunk().chunk_size:
                final_chunks.append(group)
            else:
                final_chunks.extend(self._split_text_fixed(group))
        return final_chunks

    def _split_by_structure(self, text: str) -> list[str]:
        starts = sorted({match.start() for match in STRUCTURE_SPLIT_PATTERN.finditer(text)} | {0})
        if len(starts) == 1 and starts[0] == 0:
            paragraphs = [item.strip() for item in text.split("\n\n") if item.strip()]
            if len(paragraphs) > 1:
                return paragraphs
            return [text.strip()] if text.strip() else []

        sections: list[str] = []
        for index, start in enumerate(starts):
            end = starts[index + 1] if index + 1 < len(starts) else len(text)
            piece = text[start:end].strip()
            if piece:
                sections.append(piece)
        return sections

    def _split_text_structure(self, text: str) -> list[str]:
        sections = self._split_by_structure(text)
        if not sections:
            return []

        chunks: list[str] = []
        for section in sections:
            cleaned = section.strip()
            if not cleaned:
                continue
            if len(cleaned) <= self._chunk().chunk_size:
                chunks.append(cleaned)
            else:
                chunks.extend(self._split_text_recursive(cleaned))
        return chunks

    def _embed_texts(self, texts: list[str]) -> list[list[float]]:
        if self.settings.embedding_provider == "remote":
            response = self._get_embedding_client().embeddings.create(
                model=self.settings.embedding_model,
                input=texts,
            )
            return [item.embedding for item in response.data]

        if self.settings.embedding_provider == "local":
            model = self._get_local_embedding_model()
        else:
            model = self._get_huggingface_embedding_model()
        embeddings = model.encode(texts, normalize_embeddings=True)
        return embeddings.tolist()

    def _tokenize_text(self, text: str) -> list[str]:
        lowered = text.lower()
        word_tokens = re.findall(r"[a-zA-Z0-9_]+", lowered)
        cjk_tokens = re.findall(r"[\u4e00-\u9fff]", lowered)
        return word_tokens + cjk_tokens

    def _keyword_score(self, question: str, document: str) -> float:
        query_tokens = self._tokenize_text(question)
        if not query_tokens:
            return 0.0

        doc_text = document.lower()
        score = 0.0
        for token in query_tokens:
            if token and token in doc_text:
                score += 1.0
        return score / max(len(query_tokens), 1)

    def _cosine_similarity(self, vector_a: list[float], vector_b: list[float]) -> float:
        if vector_a is None or vector_b is None:
            return 0.0
        if hasattr(vector_a, "tolist"):
            vector_a = vector_a.tolist()
        if hasattr(vector_b, "tolist"):
            vector_b = vector_b.tolist()
        if len(vector_a) == 0 or len(vector_b) == 0 or len(vector_a) != len(vector_b):
            return 0.0

        dot = sum(a * b for a, b in zip(vector_a, vector_b))
        norm_a = math.sqrt(sum(a * a for a in vector_a))
        norm_b = math.sqrt(sum(b * b for b in vector_b))
        if norm_a == 0 or norm_b == 0:
            return 0.0
        return dot / (norm_a * norm_b)

    def _keyword_retrieve(self, question: str, collection_name: str | None = None) -> list[dict[str, Any]]:
        collection = self._get_collection(collection_name)
        result = collection.get(include=["documents", "metadatas", "embeddings"])
        documents = result.get("documents", [])
        metadatas = result.get("metadatas", [])

        ranked: list[dict[str, Any]] = []
        for index, document in enumerate(documents):
            score = self._keyword_score(question, document or "")
            if score <= 0:
                continue
            metadata = metadatas[index] if index < len(metadatas) else {}
            ranked.append(
                _context_item_from_metadata(
                    metadata,
                    document or "",
                    distance=round(1 - score, 6),
                    score=score,
                )
            )

        ranked.sort(key=lambda item: item["score"], reverse=True)
        return ranked[: self.settings.top_k]

    def _keyword_retrieve_candidates(
        self,
        question: str,
        n_results: int,
        collection_name: str | None = None,
    ) -> list[dict[str, Any]]:
        collection = self._get_collection(collection_name)
        result = collection.get(include=["documents", "metadatas", "embeddings"])
        documents = result.get("documents", [])
        metadatas = result.get("metadatas", [])

        ranked: list[dict[str, Any]] = []
        for index, document in enumerate(documents):
            score = self._keyword_score(question, document or "")
            if score <= 0:
                continue
            metadata = metadatas[index] if index < len(metadatas) else {}
            ranked.append(
                _context_item_from_metadata(
                    metadata,
                    document or "",
                    distance=round(1 - score, 6),
                    score=score,
                )
            )

        ranked.sort(key=lambda item: item["score"], reverse=True)
        return ranked[:n_results]

    def _vector_retrieve(
        self,
        question: str,
        n_results: int,
        collection_name: str | None = None,
    ) -> list[dict[str, Any]]:
        collection = self._get_collection(collection_name)
        if collection.count() == 0:
            return []

        question_embedding = self._embed_texts([question])[0]
        result = collection.query(
            query_embeddings=[question_embedding],
            n_results=n_results,
            include=["documents", "metadatas", "distances", "embeddings"],
        )

        documents = result.get("documents", [[]])[0]
        metadatas = result.get("metadatas", [[]])[0]
        distances = result.get("distances", [[]])[0]
        embeddings = result.get("embeddings", [[]])[0]

        items: list[dict[str, Any]] = []
        for index, document in enumerate(documents):
            metadata = metadatas[index] if index < len(metadatas) else {}
            distance = distances[index] if index < len(distances) else None
            embedding = embeddings[index] if index < len(embeddings) else None
            items.append(
                _context_item_from_metadata(
                    metadata,
                    document or "",
                    distance=distance,
                    embedding=embedding,
                )
            )
        return items

    def _rerank_results(self, question: str, items: list[dict[str, Any]]) -> list[dict[str, Any]]:
        question_embedding = self._embed_texts([question])[0]
        ranked: list[dict[str, Any]] = []

        for item in items:
            vector_score = 0.0
            if item.get("embedding") is not None:
                vector_score = self._cosine_similarity(question_embedding, item["embedding"])
            keyword_score = self._keyword_score(question, item.get("content", ""))
            combined_score = (vector_score * 0.7) + (keyword_score * 0.3)
            ranked.append(
                {
                    **item,
                    "score": combined_score,
                    "distance": round(1 - combined_score, 6),
                }
            )

        ranked.sort(key=lambda item: item["score"], reverse=True)
        return [
            {key: value for key, value in item.items() if key != "embedding"}
            for item in ranked[: self.settings.top_k]
        ]

    def _rrf_fuse_results(
        self,
        keyword_items: list[dict[str, Any]],
        vector_items: list[dict[str, Any]],
    ) -> list[dict[str, Any]]:
        fused_scores: dict[tuple[str, str], dict[str, Any]] = {}

        for rank, item in enumerate(keyword_items, start=1):
            key = (item.get("source", "unknown"), item.get("content", ""))
            fused_scores.setdefault(key, _seed_fused_item(item))
            fused_scores[key]["score"] += 1.0 / (self.settings.rrf_k + rank)

        for rank, item in enumerate(vector_items, start=1):
            key = (item.get("source", "unknown"), item.get("content", ""))
            fused_scores.setdefault(key, _seed_fused_item(item))
            fused_scores[key]["score"] += 1.0 / (self.settings.rrf_k + rank)

        ranked = sorted(
            fused_scores.values(),
            key=lambda item: item["score"],
            reverse=True,
        )

        return [
            {
                key: value
                for key, value in {
                    **item,
                    "distance": round(1 - item["score"], 6),
                }.items()
                if key != "score"
            }
            for item in ranked[: self.settings.top_k]
        ]

    def _reset_chroma_storage(self, reason: str) -> None:
        if not self.settings.chroma_dir.exists():
            return

        backup_dir = self.settings.chroma_dir.with_name(
            f"{self.settings.chroma_dir.name}_broken_backup"
        )
        if backup_dir.exists():
            shutil.rmtree(backup_dir, ignore_errors=True)
        shutil.move(str(self.settings.chroma_dir), str(backup_dir))
        self.settings.chroma_dir.mkdir(parents=True, exist_ok=True)
        logger.warning(
            "检测到 Chroma 存储异常，已自动备份旧目录并重建。reason=%s backup=%s",
            reason,
            backup_dir,
        )

    def _get_vector_client(self) -> Any:
        if self._vector_client is not None:
            return self._vector_client

        try:
            self._vector_client = chromadb.PersistentClient(path=str(self.settings.chroma_dir))
            return self._vector_client
        except Exception as exc:
            logger.exception("初始化 Chroma PersistentClient 失败，准备重建存储目录")
            self._reset_chroma_storage(str(exc))
            self._vector_client = chromadb.PersistentClient(path=str(self.settings.chroma_dir))
            logger.info("Chroma PersistentClient 已在重建后的目录中恢复")
            return self._vector_client

    def _get_collection(self, collection_name: str | None = None) -> Any:
        name = collection_name or self.settings.collection_name
        return self._get_vector_client().get_or_create_collection(
            name=name,
            metadata={"hnsw:space": "cosine"},
        )

    def build_index(
        self,
        reset: bool = True,
        batch_size: int = 32,
        profile: KnowledgeBaseProfile | None = None,
    ) -> dict[str, int | str]:
        collection_name = profile.collection_name if profile else self.settings.collection_name
        chunk_config = profile.to_chunk_config() if profile else None
        logger.info(
            "开始读取文档并构建索引: collection=%s chunk_strategy=%s embedding_provider=%s batch_size=%s",
            collection_name,
            (chunk_config.chunk_strategy if chunk_config else self.settings.chunk_strategy),
            self.settings.embedding_provider,
            batch_size,
        )
        documents, source_summary = self.load_documents()
        if not documents:
            raise FileNotFoundError(
                "没有读取到可用数据，请先放入知识库文件，或确认 MySQL 配置已经正确填写。"
            )

        if reset:
            try:
                self._get_vector_client().delete_collection(collection_name)
            except Exception:
                pass

        collection = self._get_collection(collection_name)

        all_ids: list[str] = []
        all_chunks: list[str] = []
        all_metadatas: list[dict[str, Any]] = []

        profile_prefix = profile.id if profile else "default"
        for doc_index, document in enumerate(documents):
            chunks = self.split_text(document["content"], chunk_config=chunk_config)
            for chunk_index, chunk in enumerate(chunks):
                all_ids.append(f"{profile_prefix}-doc-{doc_index}-chunk-{chunk_index}")
                all_chunks.append(chunk)
                all_metadatas.append(
                    _chunk_metadata_from_document(
                        document,
                        chunk_index=chunk_index,
                        kb_id=profile.id if profile else "",
                        collection_name=collection_name,
                    )
                )

        if not all_chunks:
            raise ValueError("已读取到数据，但切分后没有产生可索引内容。")

        for start in range(0, len(all_chunks), batch_size):
            end = start + batch_size
            batch_chunks = all_chunks[start:end]
            batch_embeddings = self._embed_texts(batch_chunks)
            collection.upsert(
                ids=all_ids[start:end],
                documents=batch_chunks,
                metadatas=all_metadatas[start:end],
                embeddings=batch_embeddings,
            )

        logger.info(
            "索引构建写入完成: collection=%s documents=%s chunks=%s",
            collection_name,
            source_summary.get("documents"),
            len(all_chunks),
        )
        return {
            **source_summary,
            "collection_name": collection_name,
            "chunks": len(all_chunks),
        }

    def build_multi_knowledge_bases(
        self,
        profiles: list[KnowledgeBaseProfile],
        *,
        reset: bool = True,
        batch_size: int = 32,
    ) -> dict[str, Any]:
        if not profiles:
            raise ValueError("请至少选择一个知识库配置。")

        summaries: dict[str, Any] = {}
        for profile in profiles:
            summaries[profile.id] = self.build_index(reset=reset, batch_size=batch_size, profile=profile)
        summaries["profiles"] = [profile.id for profile in profiles]
        summaries["collections"] = [profile.collection_name for profile in profiles]
        return summaries

    def _retrieve_from_collection(
        self,
        question: str,
        collection_name: str,
    ) -> list[dict[str, Any]]:
        if self.settings.retrieval_method == "keyword":
            items = self._keyword_retrieve(question, collection_name=collection_name)
        elif self.settings.retrieval_method == "vector":
            vector_items = self._vector_retrieve(
                question,
                n_results=self.settings.top_k,
                collection_name=collection_name,
            )
            items = [{**item, "collection_name": collection_name} for item in vector_items]
        elif self.settings.retrieval_method == "rerank":
            vector_items = self._vector_retrieve(
                question,
                n_results=self.settings.rerank_candidate_k,
                collection_name=collection_name,
            )
            items = self._rerank_results(question, vector_items)
            for item in items:
                item["collection_name"] = collection_name
        else:
            keyword_items = self._keyword_retrieve_candidates(
                question,
                self.settings.rerank_candidate_k,
                collection_name=collection_name,
            )
            vector_items = self._vector_retrieve(
                question,
                n_results=self.settings.rerank_candidate_k,
                collection_name=collection_name,
            )
            items = self._rrf_fuse_results(keyword_items, vector_items)
            for item in items:
                item["collection_name"] = collection_name
        return items

    def retrieve_multi(
        self,
        question: str,
        profiles: list[KnowledgeBaseProfile],
        timer: StageTimer | None = None,
    ) -> list[dict[str, Any]]:
        if not profiles:
            return self.retrieve(question, timer=timer)

        if timer:
            started_at = timer.start("vector_retrieve")
            logger.info("STAGE_START | vector_retrieve | started_at=%s", started_at)

        per_collection_k = max(self.settings.top_k, self.settings.rerank_candidate_k)
        groups: list[list[dict[str, Any]]] = []
        workers = min(len(profiles), 4)

        def _fetch(profile: KnowledgeBaseProfile) -> list[dict[str, Any]]:
            items = self._retrieve_from_collection(question, profile.collection_name)
            for item in items:
                item["kb_id"] = profile.id
                item["kb_name"] = profile.name
            return items[:per_collection_k]

        if workers <= 1:
            groups = [_fetch(profile) for profile in profiles]
        else:
            with ThreadPoolExecutor(max_workers=workers) as executor:
                future_map = {executor.submit(_fetch, profile): profile for profile in profiles}
                for future in as_completed(future_map):
                    groups.append(future.result())

        items = merge_retrieval_results(groups, self.settings.top_k)
        logger.info(
            "多知识库并行检索完成: profiles=%s merged_hits=%s",
            [profile.id for profile in profiles],
            len(items),
        )

        if timer:
            record = timer.end("vector_retrieve")
            logger.info(
                "STAGE_END | %s | started_at=%s | ended_at=%s | duration_ms=%s",
                record.name,
                record.started_at,
                record.ended_at,
                record.duration_ms,
            )
        return items

    def retrieve(self, question: str, timer: StageTimer | None = None) -> list[dict[str, Any]]:
        profiles = resolve_retrieval_profiles(self.settings)
        if profiles:
            return self.retrieve_multi(question, profiles, timer=timer)

        if timer:
            started_at = timer.start("vector_retrieve")
            logger.info("STAGE_START | vector_retrieve | started_at=%s", started_at)

        logger.info("当前检索模式: %s", self.settings.retrieval_method)
        items = self._retrieve_from_collection(question, self.settings.collection_name)

        if timer:
            record = timer.end("vector_retrieve")
            logger.info(
                "STAGE_END | %s | started_at=%s | ended_at=%s | duration_ms=%s",
                record.name,
                record.started_at,
                record.ended_at,
                record.duration_ms,
            )
        return items

    def answer(self, question: str, timer: StageTimer | None = None) -> dict[str, Any]:
        if not self.settings.chat_model:
            raise ValueError("请先在配置中填写 CHAT_MODEL。")

        logger.info("开始执行 RAG 问答: question=%s", question)
        contexts = self.retrieve(question, timer=timer)
        if not contexts:
            raise RuntimeError("当前向量库为空，请先执行 ingest.py 建立索引。")

        if timer:
            started_at = timer.start("build_prompt_context")
            logger.info("STAGE_START | build_prompt_context | started_at=%s", started_at)
        context_text = "\n\n".join(
            self._format_context_block(index, item)
            for index, item in enumerate(contexts, start=1)
        )
        references = list(dict.fromkeys(item["source"] for item in contexts))
        if timer:
            record = timer.end("build_prompt_context")
            logger.info(
                "STAGE_END | %s | started_at=%s | ended_at=%s | duration_ms=%s",
                record.name,
                record.started_at,
                record.ended_at,
                record.duration_ms,
            )

        if timer:
            started_at = timer.start("llm_completion")
            logger.info("STAGE_START | llm_completion | started_at=%s", started_at)
        response = self._get_openai_client().chat.completions.create(
            model=self.settings.chat_model,
            temperature=0.2,
            messages=[
                {
                    "role": "system",
                    "content": self._resolve_chat_system_prompt(),
                },
                {
                    "role": "user",
                    "content": (
                        f"以下是检索到的知识库内容：\n{context_text}\n\n"
                        f"用户问题：{question}\n\n"
                        "请用中文给出简洁准确的答案，并尽量结合上下文内容。"
                    ),
                },
            ],
        )
        if timer:
            record = timer.end("llm_completion")
            logger.info(
                "STAGE_END | %s | started_at=%s | ended_at=%s | duration_ms=%s",
                record.name,
                record.started_at,
                record.ended_at,
                record.duration_ms,
            )

        answer = (response.choices[0].message.content or "").strip()
        logger.info("RAG 问答完成: question=%s references=%s", question, references)
        return {
            "answer": answer,
            "references": references,
            "contexts": contexts,
        }

    def close(self) -> None:
        """释放缓存的客户端与模型引用，便于 API 服务刷新实例。"""
        self._vector_client = None
        self._local_embedding_model = None
        self._openai_client = None
        self._embedding_client = None
